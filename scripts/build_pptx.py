#!/usr/bin/env python3
"""Assemble the analysis figures (slides_png/) into a PowerPoint deck.

Each slide: title, the figure scaled to fit (aspect preserved), and an editable
notes box at the bottom for your own explanations. Output goes to the results
folder (gitignored), not the repo.

Usage: python3 scripts/build_pptx.py [results_dir]
"""
import os, sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

ROOT = os.path.normpath(os.path.join(os.path.dirname(os.path.abspath(__file__)), ".."))
results_dir = sys.argv[1] if len(sys.argv) > 1 else os.path.join(ROOT, "results/McEachin_DE_paper")
png_dir = os.path.join(results_dir, "slides_png")
out_path = os.path.join(results_dir, "McEachin_nitration_figures.pptx")

# (filename in slides_png, slide title, prefilled note — edit/replace on the slide)
SLIDES = [
    ("volcano_c9ALS_vs_Control.png", "Differential nitration — c9ALS vs Control",
     "Volcano of per-site nitration (limma). Dashed line = p<0.05; right = higher in c9ALS."),
    ("volcano_ASO_vs_Control.png", "Differential nitration — ASO vs Control",
     "Does BIIB078 antisense treatment shift nitration relative to controls?"),
    ("volcano_c9ALS_vs_ASO.png", "Differential nitration — c9ALS vs ASO",
     "Treatment effect: untreated c9ALS vs antisense-treated."),
    ("heatmap_c9ALS_vs_Control.png", "Nitration heatmap — c9ALS vs Control",
     "Row-scaled nitration of the p<0.05 sites across patients, grouped by diagnosis."),
    ("heatmap_ASO_vs_Control.png", "Nitration heatmap — ASO vs Control",
     "Row-scaled nitration of the p<0.05 sites; columns grouped by group."),
    ("heatmap_c9ALS_vs_ASO.png", "Nitration heatmap — c9ALS vs ASO",
     "Row-scaled nitration of the p<0.05 sites; columns grouped by group."),
    ("PPI_c9ALS_vs_Control_pathways.png", "STRING network + pathways — c9ALS vs Control",
     "Significant-gene interaction network with enriched KEGG/Reactome/WikiPathways."),
    ("PPI_c9ALS_vs_Control_GO-BP.png", "STRING network + GO-BP — c9ALS vs Control",
     "Same network with enriched GO Biological Process terms."),
    ("PPI_ASO_vs_Control_pathways.png", "STRING network + pathways — ASO vs Control",
     "Significant-gene interaction network with enriched pathways."),
    ("PPI_ASO_vs_Control_GO-BP.png", "STRING network + GO-BP — ASO vs Control",
     "Same network with enriched GO Biological Process terms."),
    ("PPI_c9ALS_vs_ASO_pathways.png", "STRING network + pathways — c9ALS vs ASO",
     "Significant-gene interaction network with enriched pathways."),
    ("PPI_c9ALS_vs_ASO_GO-BP.png", "STRING network + GO-BP — c9ALS vs ASO",
     "Same network with enriched GO Biological Process terms."),
    ("variance_moderation.png", "limma variance moderation (empirical Bayes)",
     "Why limma differs from Welch: per-site variance is shrunk toward the prior s0^2."),
    ("method_significance_counts.png", "Significant sites by method",
     "Count of p<0.05 sites, limma vs Welch, for each contrast."),
    ("concordance_c9ALS_vs_Control.png", "limma vs Welch concordance — c9ALS vs Control",
     "Per-site p agreement; red = limma-only, blue = Welch-only calls."),
    ("concordance_ASO_vs_Control.png", "limma vs Welch concordance — ASO vs Control",
     "Per-site p agreement; red = limma-only, blue = Welch-only calls."),
    ("concordance_c9ALS_vs_ASO.png", "limma vs Welch concordance — c9ALS vs ASO",
     "Per-site p agreement; red = limma-only, blue = Welch-only calls."),
]

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
blank = prs.slide_layouts[6]

NAVY = RGBColor(0x1F, 0x3A, 0x5F)
GREY = RGBColor(0x80, 0x80, 0x80)

# --- title slide ---
s = prs.slides.add_slide(blank)
tb = s.shapes.add_textbox(Inches(0.8), Inches(2.4), SW - Inches(1.6), Inches(1.6)).text_frame
tb.word_wrap = True
tb.text = "Protein Nitration in C9orf72 ALS Spinal Cord"
tb.paragraphs[0].font.size = Pt(34); tb.paragraphs[0].font.bold = True
tb.paragraphs[0].font.color.rgb = NAVY
p = tb.add_paragraph(); p.text = "Re-analysis figures — differential nitration, enrichment/PPI, and method comparison"
p.font.size = Pt(18); p.font.color.rgb = GREY
sub = s.shapes.add_textbox(Inches(0.8), Inches(4.2), SW - Inches(1.6), Inches(0.6)).text_frame
sub.text = "Arpit Ramani · Roberts Lab, Emory"
sub.paragraphs[0].font.size = Pt(16)

# --- figure slides ---
TITLE_H = Inches(0.85)
NOTE_H = Inches(1.15)
MARGIN = Inches(0.4)
img_top = TITLE_H + Inches(0.05)
img_max_h = SH - img_top - NOTE_H - Inches(0.15)
img_max_w = SW - 2 * MARGIN

missing = []
for fname, title, note in SLIDES:
    path = os.path.join(png_dir, fname)
    if not os.path.exists(path):
        missing.append(fname); continue
    s = prs.slides.add_slide(blank)

    t = s.shapes.add_textbox(MARGIN, Inches(0.18), SW - 2 * MARGIN, TITLE_H).text_frame
    t.word_wrap = True
    t.text = title
    t.paragraphs[0].font.size = Pt(24); t.paragraphs[0].font.bold = True
    t.paragraphs[0].font.color.rgb = NAVY

    iw, ih = Image.open(path).size
    scale = min(img_max_w / iw, img_max_h / ih)
    w = int(iw * scale); h = int(ih * scale)
    left = int((SW - w) / 2); top = int(img_top + (img_max_h - h) / 2)
    s.shapes.add_picture(path, left, top, width=w, height=h)

    nb = s.shapes.add_textbox(MARGIN, SH - NOTE_H - Inches(0.1), SW - 2 * MARGIN, NOTE_H)
    nf = nb.text_frame; nf.word_wrap = True; nf.vertical_anchor = MSO_ANCHOR.TOP
    nf.text = "Notes: " + note
    nf.paragraphs[0].font.size = Pt(13)
    nf.paragraphs[0].font.color.rgb = GREY
    # faint box so the editable area is obvious on the slide
    ln = nb.line; ln.color.rgb = RGBColor(0xCC, 0xCC, 0xCC); ln.width = Pt(0.75)

prs.save(out_path)
print("wrote", out_path, "with", len(prs.slides), "slides")
if missing:
    print("missing (skipped):", ", ".join(missing))
